import shutil
import os
import re
import io
import sys
import traceback

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import pdfplumber
from reportlab.pdfgen import canvas
from PyPDF2 import PdfReader, PdfWriter
import xlrd # for .xls
from xlutils.copy import copy as xl_copy
import base64
import asyncio
import fitz  # PyMuPDF
# from configparser import ConfigParser
from presidio_analyzer import AnalyzerEngine
from presidio_analyzer.nlp_engine import SpacyNlpEngine
import spacy
import boto3
import json
import logging

log_capture_string = io.StringIO()
ch = logging.StreamHandler(log_capture_string)
ch.setLevel(logging.DEBUG)

formatter = logging.Formatter('%(asctime)s %(levelname)-8s %(message)s', datefmt='%Y-%m-%d %H:%M:%S')
ch.setFormatter(formatter)

# Get logger
logger = logging.getLogger()
logger.setLevel(logging.INFO)
logger.addHandler(ch)

# =====================================================
# ENVIRONMENT VARIABLES
# =====================================================
# secret_name = "manual_lambda_secret"
# region_name = "eu-west-1"
secret_name = os.environ['secret_name']
region_name = os.environ['region_name']


# =====================================================
# RETRIEVE SECRETS
# =====================================================
def get_secret(secret_name, region_name):
    session = boto3.session.Session()
    client = session.client("secretsmanager", region_name=region_name)
    secret = json.loads(client.get_secret_value(SecretId=secret_name)["SecretString"])
    return json.loads(secret['configPythonSecrets'])

config = get_secret(secret_name, region_name)

BEDROCK_REGION = config['awsRegion']
BEDROCK_MODEL_ID = config['bedrock']['llm']

# # Configure Bedrock details via environment variables
# BEDROCK_REGION = "eu-west-1"
# BEDROCK_MODEL_ID = "arn:aws:bedrock:eu-west-1:572431880657:inference-profile/eu.anthropic.claude-3-7-sonnet-20250219-v1:0"

IMAGE_SAVE_DIR_D = "docx_images"
IMAGE_SAVE_DIR_E = "excel_images"
IMAGE_SAVE_FOLDER_PPT = "ppt_images"
IMAGE_SAVE_DIR_P = "pdf_images"

async def bedrock_invoke_model_async(model_id: str, input_payload: dict, region_name: str = None) -> dict:
    """Invoke Bedrock model in a thread to avoid blocking the event loop.

    This returns parsed JSON if possible, otherwise raw text under 'text'.
    """
    region = region_name or BEDROCK_REGION

    def _invoke():
        client = boto3.client("bedrock-runtime", region_name=region)
        response = client.invoke_model(
            modelId=model_id,
            # contentType="application/json",
            # accept="application/json",
            body=json.dumps(input_payload),
        )
        # response['body'] is a StreamingBody
        body_bytes = response["body"].read()
        try:
            return json.loads(body_bytes.decode("utf-8"))
        except Exception:
            # return raw text fallback
            try:
                return {"text": body_bytes.decode("utf-8")}
            except Exception:
                return {"text": ""}

    return await asyncio.to_thread(_invoke)


def _messages_to_prompt(messages) -> str:
    """Flatten a messages structure (similar to OpenAI chat format) to a single prompt string.

    Supports messages where content is either a string or a list of dicts containing text/image_url objects.
    """
    parts = []
    if isinstance(messages, str):
        return messages
    for m in messages:
        role = m.get("role", "user")
        content = m.get("content")
        if isinstance(content, str):
            parts.append(f"[{role.upper()}] {content}")
        elif isinstance(content, list):
            # handle list of items like {type: 'text', text: '...'} or image_url
            for item in content:
                if item.get("type") == "text":
                    text = item.get("text", "")
                    parts.append(f"[{role.upper()}] {text}")
                elif item.get("type") == "image_url":
                    url_obj = item.get("image_url", {})
                    url = url_obj.get("url")
                    # include a short marker for image; avoid sending huge base64 blobs inline repeatedly
                    if url and url.startswith("data:image"):
                        # keep a truncated preview of the base64 to help model if needed
                        preview = url[:2000]
                        parts.append(f"[{role.upper()}] IMAGE_DATA: {preview}")
                    else:
                        parts.append(f"[{role.upper()}] IMAGE_URL: {url}")
        else:
            parts.append(f"[{role.upper()}] {str(content)}")
    return "\n".join(parts)


async def bedrock_chat(model_id: str, messages, temperature: float = 0.0) -> str:
    """High-level chat helper that converts messages into a prompt and calls Bedrock.

    Returns the text response from the model.
    """
    # If 'messages' is already a list of message dicts (role/content), pass it through.
    #print(f"In bedrock_chat with messages: {messages}")
    if isinstance(messages, list):
        payload_messages = []
        for m in messages:
            # ensure content structure: a list of items with type/text or image_url
            content = m.get("content")
            if isinstance(content, str):
                payload_messages.append({"role": m.get("role", "user"), "content": [{"type": "text", "text": content}]})
            else:
                payload_messages.append({"role": m.get("role", "user"), "content": content})
        payload = {"messages": payload_messages, "temperature": temperature, "max_tokens": 800, "anthropic_version": "bedrock-2023-05-31",}
    else:
        # fallback: flatten into a single user message
        prompt = _messages_to_prompt(messages)
        payload = {"messages": [{"role": "user", "content": [{"type": "text", "text": prompt}]}], "temperature": temperature, "max_tokens": 800, "anthropic_version": "bedrock-2023-05-31",}

    result = await bedrock_invoke_model_async(model_id or BEDROCK_MODEL_ID, payload)
    #print(f"Result from Bedrock: {result}")
    # try common response shapes
    # 1) {'output': '...'} or {'text': '...'} or {'outputs': [{'content': '...'}]}
    if isinstance(result, dict):
        # if "text" in result and isinstance(result["text"], str):
        #     return result["text"].strip()
        # if "output" in result and isinstance(result["output"], str):
        #     return result["output"].strip()
        # if "outputs" in result and isinstance(result["outputs"], list) and result["outputs"]:
        #     first = result["outputs"][0]
        #     if isinstance(first, dict):
        #         # try several keys
        #         for k in ("content", "text", "body", "outputText"):
        #             if k in first:
        #                 return first[k].strip() if isinstance(first[k], str) else json.dumps(first[k])
        #     elif isinstance(first, str):
        #         return first.strip()
        return result["content"][0]["text"].strip() if "content" in result and isinstance(result["content"], list) and result["content"] else ""
    # fallback to empty string
    return ""

# os.makedirs(IMAGE_SAVE_DIR_D, exist_ok=True)

def _invoke_bedrock_llm(region_name, prompt, llm_model, max_tokens=800, temperature=0):
    """Invoke Bedrock model synchronously."""
    client = boto3.client("bedrock-runtime", region_name=region_name)

    body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": max_tokens,
        "temperature": temperature,
        "messages": [
            {"role": "user", "content": [{"type": "text", "text": prompt}]}
        ],
    }

    response = client.invoke_model(modelId=llm_model, body=json.dumps(body))
    model_response = json.loads(response["body"].read())

    return model_response["content"][0]["text"].strip()


def _describe_image_sync(image_path, region_name, llm_model):
    """Describe a single image synchronously using Bedrock multimodal model."""
    with open(image_path, "rb") as img_file:
        b64_img = base64.b64encode(img_file.read()).decode("utf-8")

    prompt = f"Describe this image (file: {os.path.basename(image_path)})."

    body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 300,
        "temperature": 0,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64_img}}
                ]
            }
        ]
    }

    client = boto3.client("bedrock-runtime", region_name=region_name)
    response = client.invoke_model(modelId=llm_model, body=json.dumps(body))
    model_response = json.loads(response["body"].read())

    return f"{os.path.basename(image_path)}: {model_response['content'][0]['text'].strip()}"


async def describe_images_async(image_paths, region_name, llm_model):
    """Run image descriptions concurrently by delegating to sync worker in threads."""
    tasks = []
    for img_path in image_paths:
        tasks.append(asyncio.to_thread(_describe_image_sync, img_path, region_name, llm_model))
    return await asyncio.gather(*tasks, return_exceptions=True)

# Regex patterns

EMAIL_PATTERN = re.compile(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", re.IGNORECASE)
PHONE_PATTERN = re.compile(r"\b\d{10}\b|\b\d{5}[ -]?\d{5}\b")
CREDIT_CARD_PATTERN = re.compile(r"\b(?:\d{4}[ -]?){3}\d{4}\b")  # handles 1234 5678 9012 3456 pattern
BANK_ACC_PATTERN = re.compile(r"\b\d{11,20}\b")

# Masking helper

class LoadedSpacyNlpEngine(SpacyNlpEngine):
    def __init__(self, loaded_spacy_model):
        super().__init__()
        self.nlp = {"en": loaded_spacy_model}

# Load a model a-priori
nlp = spacy.load("./en_core_web_md-3.8.0")
 
# Pass the loaded model to the new LoadedSpacyNlpEngine
loaded_nlp_engine = LoadedSpacyNlpEngine(loaded_spacy_model = nlp)


analyzer = AnalyzerEngine(nlp_engine = loaded_nlp_engine)


def mask_pii(text: str) -> str:
    if not text:
        return text

    # Mask Names using Presidio
    results = analyzer.analyze(text=text, entities=["PERSON"], language="en")
    #Sort by start index descending to avoid shifting issues
    results = sorted(results, key=lambda x: x.start, reverse=True)
    for r in results:
        start, end = r.start, r.end
        text = text[:start] + "[NAME_REDACTED]" + text[end:]

    # Step 2: Mask Emails using regex
    email_pattern = r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"
    text = re.sub(email_pattern, "[EMAIL_REDACTED]", text)

    # Step 3: Mask Phone Numbers only if + followed by >7 digits
    phone_pattern = r"\+\d{8,}[\d\s\-\(\)]*"
    text = re.sub(phone_pattern, "[PHONE_REDACTED]", text)

    return text


def mask_text(text, company_patterns=None):
    if not text:
        return text
    text = CREDIT_CARD_PATTERN.sub("[CREDIT_CARD_REDACTED]", text)
    text = BANK_ACC_PATTERN.sub("[BANK_DETAILS_REDACTED]", text)
    if company_patterns:
        for cp in company_patterns:
            text = cp.sub("[COMPANY_REDACTED]", text)
    text_final = mask_pii(text)
    return text_final


def _apply_mask_to_element(element, company_patterns):
    """
    Applies the mask_text function's logic to an element (Paragraph)
    by clearing its runs and inserting a new run with the masked text.
    This prevents the complete destruction of the parent element's XML,
    which helps preserve non-text elements like images and shapes.
    """
    full_text = element.text
    if not full_text:
        return
    masked_text = mask_text(full_text, company_patterns)
    if full_text != masked_text:
        element._element.clear_content()
        element.add_run(masked_text)

def mask_docx(file_path, company_names=None):
    doc = Document(file_path)
    company_patterns = [re.compile(rf"\b{re.escape(name)}\b", re.IGNORECASE) for name in (company_names or [])]

    for p in doc.paragraphs:
        _apply_mask_to_element(p, company_patterns)
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _apply_mask_to_element(p, company_patterns)
    return doc


# PPTX Masking

def mask_pptx(file_path, company_names=None):
    prs = Presentation(file_path)
    company_patterns = [re.compile(rf"\b{re.escape(name)}\b", re.IGNORECASE) for name in (company_names or [])]

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.text = mask_text(run.text, company_patterns)

    return prs

def mask_excel(file_path, company_names=None):
    company_patterns = [re.compile(rf"\b{re.escape(name)}\b", re.IGNORECASE) for name in (company_names or [])]
    ext = os.path.splitext(file_path)[1].lower()
    # --- Handle .xlsx files ---
    if ext == ".xlsx":
        wb = load_workbook(file_path)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        cell.value = mask_text(cell.value, company_patterns)
        return wb
    # --- Handle .xls files ---
    elif ext == ".xls":
        rb = xlrd.open_workbook(file_path, formatting_info=False)
        wb = xl_copy(rb)
        sheet_count = rb.nsheets
        for i in range(sheet_count):
            r_sheet = rb.sheet_by_index(i)
            w_sheet = wb.get_sheet(i)
            for row_idx in range(r_sheet.nrows):
                for col_idx in range(r_sheet.ncols):
                    cell_value = r_sheet.cell_value(row_idx, col_idx)
                    if isinstance(cell_value, str):
                        masked_value = mask_text(cell_value, company_patterns)
                        w_sheet.write(row_idx, col_idx, masked_value)
        return wb

# PDF Masking

def mask_pdf(file_path, company_names=None):
 
    company_patterns = [
        re.compile(rf"\b{re.escape(name)}\b", re.IGNORECASE)
        for name in (company_names or [])
    ]
    
    all_text = ""
    pdf_writer = PdfWriter()

 
    with pdfplumber.open(file_path) as pdf:
        reader = PdfReader(file_path)
        for page_num, page in enumerate(pdf.pages):
            orig_page = reader.pages[page_num]
            width, height = page.width, page.height
 
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=(width, height))
 
            # Mask text content
            text = page.extract_text() or ""
            lines = text.splitlines()
            y = height - 20
            masked_page_text = ""
            for line in lines:
                masked_line = mask_text(line, company_patterns)
                masked_page_text += masked_line + "\n"
                can.setFont("Helvetica", 12)
                can.drawString(20, y, masked_line)
                y -= 14
            all_text += masked_page_text
            
            can.save()
            packet.seek(0)
            new_pdf = PdfReader(packet)
 
            # Add masked layer (without duplication)
            if len(new_pdf.pages) > 0:
                pdf_writer.add_page(new_pdf.pages[0])
            else:
                pdf_writer.add_page(orig_page)  # fallback if no text was extracted
    
    return all_text
# ========================================================================================

# DOCUMENT PROCESSING Docx


def chunk_text(text, max_chars=9000):
    """Split long text into chunks for summarization."""
    chunks = []
    while len(text) > max_chars:
        split_at = text[:max_chars].rfind(".")
        if split_at == -1:
            split_at = max_chars
        chunks.append(text[:split_at + 1].strip())
        text = text[split_at + 1:]
    if text.strip():
        chunks.append(text.strip())
    return chunks


def extract_docx_content(doc):
    """Extract text and images from a .docx file."""

    all_text = ""
    image_paths = []

    for para in doc.paragraphs:
        all_text += para.text + "\n"

    rels = doc.part.rels
    for rel in rels.values():
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            image_filename = os.path.basename(rel.target_ref)
            image_path = os.path.join(IMAGE_SAVE_DIR_D, image_filename)
            with open(image_path, "wb") as f:
                f.write(image_data)
            image_paths.append(image_path)

    return all_text.strip(), image_paths


async def describe_images(image_paths):
    # Delegate to the threaded synchronous describer which uses the working
    # Bedrock invoke pattern (_describe_image_sync).
    return await describe_images_async(image_paths, BEDROCK_REGION, BEDROCK_MODEL_ID)

async def summarize_text_async(prompt, idx):
    messages = [
        {"role": "assistant", "content": "You summarize text clearly and concisely."},
        {"role": "user", "content": prompt},
    ]
    response_text = await bedrock_chat(BEDROCK_MODEL_ID, messages, temperature=0.3)
    return response_text


async def summarize_chunks(chunks):
    tasks = []
    for idx, chunk in enumerate(chunks, start=1):
        prompt = f"""Summarize part {idx} of the document:\n\n{chunk}
        Summarize the following  part {idx} of the document:\n\n{chunk}
        By strictly following these rules:
        A. Summary Content Requirements
        The summary must include the following on generating the summary(Don't add ** for these bullet points)
        1. Describe what the article talks about in as single line having not more than 20 words.
        2. Summarize the main components of the application described in the article by providing its name and its functionality.
        3. Summarize how the applications are integrated by providing its names, CI and its integration details.
        4. Describe potential weaknesses or failure points or symptoms which is described in article or can be inferred from components or integrations.

        B. Strict Rules
        - If the contents to be summarized is process related and not technical such as about an application or integrations, then don't summarize. Provide output as "NIL".
        - If there is no content related to components or its integrations or its points of failure of an application, Do not provide any detailsfor those sections.
        - Do not hallucinate or add any information not found in the original text.
        - Do not include any text or commentary beyond the summary.
        - Keep the tone factual and consistent with the source.



        NO NEED TO ADD THE HEADINGS OR ANY OTHER HEADERS FOR THE SUMMARY
        GIVE A PARAGRAGH SUMMARY OUTPUT
      """
        tasks.append(summarize_text_async(prompt, idx))
    return await asyncio.gather(*tasks)


# -------------------- STEP 4: Combine into Final Summary --------------------

async def summarize_docx_with_images(docx_path):
    text, image_paths = extract_docx_content(docx_path)
    # print("Masked text in docx: {text}")
    #print(f"Extracted text: {text}")
    text_chunks = chunk_text(text)
    chunk_summaries, image_descriptions = await asyncio.gather(
        summarize_chunks(text_chunks),
        describe_images(image_paths)
    )
    #print(f"Chunk summaries: {chunk_summaries}")
    combined_input = f"""
Below are the partial summaries of the document:
{chunk_summaries}

And here are the image descriptions extracted from the PDF:
{image_descriptions}

By strictly following these rules:
A. Summary Content Requirements
The summary must include the following on generating the summary(Don't add ** for these bullet points)
1. Describe what the article talks about in as single line having not more than 20 words.
2. Summarize the main components of the application described in the article by providing its name and its functionality.
3. Summarize how the applications are integrated by providing its names, CI and its integration details.
4. Describe potential weaknesses or failure points or symptoms which is described in article or can be inferred from components or integrations.

B. Strict Rules
- If the contents to be summarized is process related and not technical such as about an application or integrations, then don't summarize. Provide output as "NIL".
- If there is no content related to components or its integrations or its points of failure of an application, Do not provide any details for those sections.
- Do not hallucinate or add any information not found in the original text.
- Do not include any text or commentary beyond the summary.
- Keep the tone factual and consistent with the source.

Also follow to give your summary output by following these rules:
Summary Length Rules
Maximum summary length can be 300 words
Minimum summary length can be 20 words or less

NO NEED TO ADD THE HEADINGS OR ANY OTHER HEADERS FOR THE SUMMARY
GIVE A PARAGRAGH SUMMARY OUTPUT

"""


    messages = [
        {"role": "assistant", "content": "You are an expert summarizer for long documents with visuals."},
        {"role": "user", "content": combined_input},
    ]
    final_text = await bedrock_chat(BEDROCK_MODEL_ID, messages, temperature=0.4)

    return {
        "final_summary": final_text,
        "chunk_summaries": chunk_summaries,
        "image_descriptions": image_descriptions,
        "image_paths": image_paths
    }
IMAGES_SAVE_DOC = "doc_images"


# -------------------- PPT --------------------------------------

# os.makedirs(IMAGE_SAVE_FOLDER_PPT, exist_ok=True)
# ==========EXTRACT PPT CONTENT ==========
def extract_ppt_content(prs, image_folder=IMAGE_SAVE_FOLDER_PPT):
    # prs = Presentation(ppt_path)
    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        text_content = ""
        images_b64 = []
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                text_content += shape.text + "\n"
            # Extract image
            if shape.shape_type == 13:  # 13 = PICTURE
                image = shape.image
                img_bytes = image.blob
                img_path = os.path.join(image_folder, f"slide_{i}_img_{len(images_b64)+1}.png")
                # Save image locally
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                images_b64.append(img_path)
        slides_data.append({
            "slide_number": i,
            "text": text_content.strip(),
            "images": images_b64
        })
    return slides_data

# ==========HELPER: CHUNK TEXT ==========
def chunk_tex_ppt(text, max_tokens=1500):
   words = text.split()
   chunk_size = max_tokens * 4  # roughly 4 tokens per word
   for i in range(0, len(words), chunk_size):
       yield " ".join(words[i:i + chunk_size])

# ==========SUMMARIZE TEXT CHUNK ==========
async def summarize_text_async_ppt(prompt, idx, client, deployment_name):
    #print(f"Summarizing PPT chunk :{prompt}")
    system_msg = """
            By strictly following these rules:
            A. Summary Content Requirements
            The summary must include the following on generating the summary(Don't add ** for these bullet points)
            1. Describe what the article talks about in as single line having not more than 20 words.
            2. Summarize the main components of the application described in the article by providing its name and its functionality.
            3. Summarize how the applications are integrated by providing its names, CI and its integration details.
            4. Describe potential weaknesses or failure points or symptoms which is described in article or can be inferred from components or integrations.

            B. Strict Rules
                    - If the contents to be summarized is process related and not technical such as about an application or integrations, then don't summarize. Provide output as "NIL".
                    - If there is no content related to components or its integrations or its points of failure of an application, Do not provide anydetails for those sections.
                    - Do not hallucinate or add any information not found in the original text.
                    - Do not include any text or commentary beyond the summary.
                    - Keep the tone factual and consistent with the source.

            NO NEED TO ADD THE HEADINGS OR ANY OTHER HEADERS FOR THE SUMMARY
            GIVE A PARAGRAGH SUMMARY OUTPUT
                    """
    messages = [
        {"role": "assistant", "content": system_msg},
        {"role": "user", "content": prompt}
    ]
    response_text = await bedrock_chat(BEDROCK_MODEL_ID, messages, temperature=0.3)
    #print(f"PPT Chunk {idx} Summary: {response_text}")
    return response_text

# ==========MAIN PIPELINE ==========
async def process_ppt(prs):
    slides_data = extract_ppt_content(prs)
    all_text = ""
    image_tasks = []
    # Collect all image paths and run the threaded describer
    all_image_paths = []
    for slide in slides_data:
        for img_path in slide["images"]:
            all_image_paths.append(img_path)
    image_descriptions = await describe_images_async(all_image_paths, BEDROCK_REGION, BEDROCK_MODEL_ID) if all_image_paths else []
    for slide in slides_data:
        all_text += f"Slide {slide['slide_number']}:\n{slide['text']}\n\n"
    #print(f"Extracted PPT text: {len(all_text)}")
    #print(f"Extracted PPT images: {len(all_image_paths)} images")
    chunks = list(chunk_tex_ppt(all_text, max_tokens=1500))
    #print(f"Chunked PPT text into {len(chunks)} parts for summarization.")
    #print(chunks)
    text_tasks = [summarize_text_async_ppt(chunk, idx, None, None) for idx, chunk in enumerate(chunks)]
    summaries = await asyncio.gather(*text_tasks)
    #print(f"Summaries of PPT chunks: {summaries}")
    final_prompt = (
        "Here are summaries of PowerPoint slide texts and image descriptions. "
        f"Text Summaries:\n{summaries}\n\n"
        f"Image Descriptions:\n{image_descriptions}"
        """\n By strictly following these rules:
        By strictly following these rules:
        A. Summary Content Requirements
        The summary must include the following on generating the summary(Don't add ** for these bullet points)
        1. Describe what the article talks about in as single line having not more than 20 words.
        2. Summarize the main components of the application described in the article by providing its name and its functionality.
        3. Summarize how the applications are integrated by providing its names, CI and its integration details.
        4. Describe potential weaknesses or failure points or symptoms which is described in article or can be inferred from components or integrations.

        B. Strict Rules
                - If the contents to be summarized is process related and not technical such as about an application or integrations, then don't summarize. Provide output as "NIL".
                - If there is no content related to components or its integrations or its points of failure of an application, Do not provide anydetails for those sections.
                - Do not hallucinate or add any information not found in the original text.
                - Do not include any text or commentary beyond the summary.
                - Keep the tone factual and consistent with the source.

        Also follow to give your summary output by following these rules:
        Summary Length Rules
        Maximum summary length can be 300 words
        Minimum summary length can be 20 words or less

        NO NEED TO ADD THE HEADINGS OR ANY OTHER HEADERS FOR THE SUMMARY
        GIVE A PARAGRAGH SUMMARY OUTPUT

    """

    )
    final_summary = await summarize_text_async_ppt(final_prompt, 0, None, None)
    return final_summary,image_descriptions


# ------------------- EXCEL -----------------------------------------------------------

# os.makedirs(IMAGE_SAVE_DIR_E, exist_ok=True)

def chunk_text(text, max_chars=9000):
    chunks = []
    while len(text) > max_chars:
        split_at = text[:max_chars].rfind(".")
        if split_at == -1:
            split_at = max_chars
        chunks.append(text[:split_at + 1].strip())
        text = text[split_at + 1:]
    if text.strip():
        chunks.append(text.strip())
    return chunks

def extract_excel_content(wb): # works for .xlsx
    """Extract text and images from Excel file."""
    all_text = ""
    image_paths = []

    # wb = openpyxl.load_workbook(excel_path, data_only=True)
    for sheet in wb.worksheets:
        all_text += f"\nSheet: {sheet.title}\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            all_text += row_text + "\n"

        # Extract images
        for idx, image in enumerate(sheet._images, start=1):
            img_bytes = image.ref  # openpyxl stores reference
            img_obj = image._data()
            image_filename = f"{sheet.title}_img{idx}.png"
            image_path = os.path.join(IMAGE_SAVE_DIR_E, image_filename)
            with open(image_path, "wb") as f:
                f.write(img_obj)
            image_paths.append(image_path)

    return all_text.strip(), image_paths

# -------------------- STEP 2: Describe Images --------------------
async def describe_images(image_paths):
    # Use the threaded synchronous describer that follows the working Bedrock payload pattern
    return await describe_images_async(image_paths, BEDROCK_REGION, BEDROCK_MODEL_ID)

# -------------------- STEP 4: Combine into Final Summary --------------------
async def summarize_excel_with_images(excel):
    # 'excel' is a workbook object returned from mask_excel; extract contents
    text, image_paths = extract_excel_content(excel)
    #print(f"Extracted text: {text}")
    #print(f"Extracted images: {len(image_paths)}")
    text_chunks = chunk_text(text)
    chunk_summaries, image_descriptions = await asyncio.gather(
        summarize_chunks(text_chunks),
        describe_images(image_paths)
    )

    combined_input = f"""
Below are the partial summaries of the document:
{chunk_summaries}

And here are the image descriptions extracted from the PDF:
{image_descriptions}

By strictly following these rules:
A. Summary Content Requirements
The summary must include the following on generating the summary(Don't add ** for these bullet points)
1. Describe what the article talks about in as single line having not more than 20 words.
2. Summarize the main components of the application described in the article by providing its name and its functionality.
3. Summarize how the applications are integrated by providing its names, CI and its integration details.
4. Describe potential weaknesses or failure points or symptoms which is described in article or can be inferred from components or integrations.

B. Strict Rules
        - If the contents to be summarized is process related and not technical such as about an application or integrations, then don't summarize. Provide output as "NIL".
        - If there is no content related to components or its integrations or its points of failure of an application, Do not provide any detailsfor those sections.
        - Do not hallucinate or add any information not found in the original text.
        - Do not include any text or commentary beyond the summary.
        - Keep the tone factual and consistent with the source.

Also follow to give your summary output by following these rules:
Summary Length Rules
Maximum summary length can be 300 words
Minimum summary length can be 20 words or less

NO NEED TO ADD THE HEADINGS OR ANY OTHER HEADERS FOR THE SUMMARY
GIVE A PARAGRAGH SUMMARY OUTPUT
          """

    messages = [
        {"role": "assistant", "content": "You are an expert summarizer for long documents with visuals."},
        {"role": "user", "content": combined_input},
    ]
    final_text = await bedrock_chat(BEDROCK_MODEL_ID, messages, temperature=0.4)

    return {
        "final_summary": final_text,
        "chunk_summaries": chunk_summaries,
        "image_descriptions": image_descriptions,
        "image_paths": image_paths
    }



# -------------------- PDF ---------------------------------

# Check for previous code
def chunk_text(text, max_chars=9000):
    """Split long text into chunks for summarization."""
    chunks = []
    while len(text) > max_chars:
        split_at = text[:max_chars].rfind(".")
        if split_at == -1:
            split_at = max_chars
        chunks.append(text[:split_at+1].strip())
        text = text[split_at+1:]
    if text.strip():
        chunks.append(text.strip())
    return chunks

# --- STEP 1: Extract PDF text & images and save images locally ---

def extract_pdf_content(pdf_path):
    doc = fitz.open(pdf_path)
    # doc = fitz.open(pdf)
    image_paths = []

    for page_num, page in enumerate(doc, start=1):
        # all_text += page.get_text("text") + "\n"
        for img_index, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            # if image_bytes:
                #print("IMAGE FOUND")
            # Construct a filename
            image_filename = f"page{page_num}_img{img_index}.png"
            image_path = os.path.join(IMAGE_SAVE_DIR_P, image_filename)
            # Save image
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            image_paths.append(image_path)
    return image_paths

async def summarize_pdf_with_images(pdf, text):
    image_paths = extract_pdf_content(pdf)
    text_chunks = chunk_text(text.strip())
    chunk_summaries, image_descriptions = await asyncio.gather(
        summarize_chunks(text_chunks),
        # describe_images(image_paths)
        describe_images_async(image_paths, BEDROCK_REGION, BEDROCK_MODEL_ID)
    )
    #image_descriptions = await describe_images_async(all_image_paths, BEDROCK_REGION, BEDROCK_MODEL_ID) if all_image_paths else []
    combined_input = f"""
Below are the partial summaries of the document:

{text_chunks}

And here are the image descriptions extracted from the PDF:

{image_descriptions}

By strictly following these rules:
A. Summary Content Requirements
The summary must include the following on generating the summary(Don't add ** for these bullet points)
1. Describe what the article talks about in as single line having not more than 20 words.
2. Summarize the main components of the application described in the article by providing its name and its functionality.
3. Summarize how the applications are integrated by providing its names, CI and its integration details.
4. Describe potential weaknesses or failure points or symptoms which is described in article or can be inferred from components or integrations.

B. Strict Rules
        - If the contents to be summarized is process related and not technical such as about an application or integrations, then don't summarize. Provide output as "NIL".
        - If there is no content related to components or its integrations or its points of failure of an application, Do not provide any detailsfor those sections.
        - Do not hallucinate or add any information not found in the original text.
        - Do not include any text or commentary beyond the summary.
        - Keep the tone factual and consistent with the source.

Also follow to give your summary output by following these rules:
Summary Length Rules
Maximum summary length can be 300 words
Minimum summary length can be 20 words or less

NO NEED TO ADD THE HEADINGS OR ANY OTHER HEADERS FOR THE SUMMARY
GIVE A PARAGRAGH SUMMARY OUTPUT
    """
    messages = [
        {"role": "assistant", "content": "You are an expert summarizer for long documents with visuals."},
        {"role": "user", "content": combined_input},
    ]
    final_text = await bedrock_chat(BEDROCK_MODEL_ID, messages, temperature=0.4)
    return {
        "final_summary": final_text,
        "chunk_summaries": chunk_summaries,
        "image_descriptions": image_descriptions,
        "image_paths": image_paths
    }


# -------------------- MAIN EXECUTION --------------------

def sharepoint_helper(file_path):
    files = [file_path]
    company_names = ["Easy Jet"] #company name
    try:
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext == ".docx":
                os.makedirs(IMAGE_SAVE_DIR_D, exist_ok=True)
                logger.info(f"Processing DOCX file: {file}")
                docx= mask_docx(file, company_names)
                result = asyncio.run(summarize_docx_with_images(docx))
                #print(f"Final Summary: {result['final_summary']}")

                shutil.rmtree(IMAGE_SAVE_DIR_D)
                return result['final_summary']

            elif ext in [".pptx"]:
                os.makedirs(IMAGE_SAVE_FOLDER_PPT, exist_ok=True)
                ppt= mask_pptx(file, company_names)
                logger.info(f"Processing PPTX file: {file}")
                summary,img_description = asyncio.run(process_ppt(ppt))
                #print(f"Final Summary: {summary}")

                shutil.rmtree(IMAGE_SAVE_FOLDER_PPT)
                return summary

            elif ext in [".xlsx"]:
                os.makedirs(IMAGE_SAVE_DIR_E, exist_ok=True)
                logger.info(f"Processing Excel file: {file}")
                excel= mask_excel(file, company_names)
                result = asyncio.run(summarize_excel_with_images(excel))
                #print(f"Final Summary: {result['final_summary']}")

                shutil.rmtree(IMAGE_SAVE_DIR_E)
                return result['final_summary']

            elif ext == ".pdf":
                os.makedirs(IMAGE_SAVE_DIR_P, exist_ok=True)
                logger.info(f"Processing PDF file: {file}")
                pdf_text = mask_pdf(file, company_names)
                result = asyncio.run(summarize_pdf_with_images(file, pdf_text))

                #print(f"Final Summary: {result['final_summary']}")
                shutil.rmtree(IMAGE_SAVE_DIR_P)
                return result['final_summary']

            else:
                logger.info(f"Unsupported file type: {ext}")
                return None
    except Exception as e:
        logger.error(f"Error processing file {file}: {traceback.format_exc()}")
